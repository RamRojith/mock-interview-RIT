import hashlib
import logging
import uuid

from django.conf import settings
from django.db import transaction
from django.utils import timezone

from mock_interview.models import DocumentChunk, UploadedDocument
from mock_interview.rag.chunker import DocumentChunker
from mock_interview.rag.embeddings import EmbeddingService
from mock_interview.rag.vectorstore import QdrantVectorStore

logger = logging.getLogger(__name__)


class DocumentServiceError(Exception):
    pass


class DocumentService:
    """Handles faculty document upload, extraction, chunking, and embedding."""

    def __init__(self):
        self._chunker = None
        self._embeddings = None
        self._vectorstore = None

    def _ensure_services(self):
        """Lazily initialize RAG services with individual error handling."""
        if self._chunker is None:
            try:
                self._chunker = DocumentChunker()
            except Exception as exc:
                logger.error("Failed to initialize document chunker: %s", exc)
                raise DocumentServiceError(
                    "Document processing configuration error. Please contact administrator."
                ) from exc

        if self._embeddings is None:
            try:
                self._embeddings = EmbeddingService()
            except Exception as exc:
                logger.error("Failed to initialize embedding model: %s", exc)
                raise DocumentServiceError(
                    "Embedding model unavailable. Ensure the BAAI/bge-m3 model is downloaded and accessible."
                ) from exc

        if self._vectorstore is None:
            try:
                self._vectorstore = QdrantVectorStore()
            except Exception as exc:
                config = settings.RAG_CONFIG
                url = config.get("QDRANT_URL", "http://127.0.0.1:6333")
                logger.error("Failed to connect to vector database at %s: %s", url, exc)
                raise DocumentServiceError(
                    f"Vector database unavailable at {url}. Please ensure Qdrant is running."
                ) from exc

    def validate_document(self, uploaded_file) -> dict:
        """Validate file type and size before processing."""
        config = settings.RAG_CONFIG
        max_bytes = config["MAX_DOCUMENT_BYTES"]
        allowed_types = config["ALLOWED_DOCUMENT_TYPES"]

        if uploaded_file.size > max_bytes:
            raise DocumentServiceError(
                f"File too large. Maximum size is {max_bytes // (1024 * 1024)} MB."
            )

        content_type = getattr(uploaded_file, "content_type", "")
        if content_type and content_type not in allowed_types:
            raise DocumentServiceError(
                f"Unsupported file type: {content_type}. "
                "Allowed: PDF, PPTX, DOCX, TXT."
            )

        return {
            "size": uploaded_file.size,
            "content_type": content_type,
        }

    def extract_text(self, uploaded_file, mime_type: str) -> str:
        """Extract text from PDF, PPTX, DOCX, or TXT files."""
        if "pdf" in mime_type:
            return self._extract_pdf(uploaded_file)
        elif "presentation" in mime_type or mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return self._extract_pptx(uploaded_file)
        elif "wordprocessingml" in mime_type or mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return self._extract_docx(uploaded_file)
        elif "text" in mime_type:
            return uploaded_file.read().decode("utf-8", errors="replace")
        else:
            raise DocumentServiceError(f"Cannot extract text from: {mime_type}")

    def _extract_pdf(self, uploaded_file) -> str:
        from pypdf import PdfReader
        import io

        uploaded_file.seek(0)
        reader = PdfReader(io.BytesIO(uploaded_file.read()))
        text_parts = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
        return "\n\n".join(text_parts)

    def _extract_pptx(self, uploaded_file) -> str:
        from pptx import Presentation
        import io

        uploaded_file.seek(0)
        prs = Presentation(io.BytesIO(uploaded_file.read()))
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text.append(text)
            if slide_text:
                text_parts.append(
                    f"[Slide {slide_num}]\n" + "\n".join(slide_text)
                )
        return "\n\n".join(text_parts)

    def _extract_docx(self, uploaded_file) -> str:
        from docx import Document
        import io

        uploaded_file.seek(0)
        doc = Document(io.BytesIO(uploaded_file.read()))
        text_parts = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                text_parts.append(text)
        return "\n".join(text_parts)

    def _compute_sha256(self, uploaded_file) -> str:
        uploaded_file.seek(0)
        digest = hashlib.sha256()
        for chunk in uploaded_file.chunks():
            digest.update(chunk)
        uploaded_file.seek(0)
        return digest.hexdigest()

    def process_document(
        self,
        uploaded_file,
        faculty_employee_id: str,
        faculty_name: str,
        subject_code: str,
        chapter: str = "",
    ) -> UploadedDocument:
        """Full pipeline: validate, extract, chunk, embed, store."""
        validation = self.validate_document(uploaded_file)
        sha256 = self._compute_sha256(uploaded_file)
        mime_type = validation["content_type"]

        with transaction.atomic():
            document = UploadedDocument.objects.create(
                faculty_employee_id=faculty_employee_id,
                faculty_name=faculty_name,
                file=uploaded_file,
                original_filename=uploaded_file.name,
                mime_type=mime_type,
                file_size=validation["size"],
                sha256=sha256,
                subject_code=subject_code,
                chapter=chapter,
                status="processing",
            )

        try:
            self._ensure_services()

            extracted_text = self.extract_text(uploaded_file, mime_type)
            document.extracted_text = extracted_text

            chunks = self._chunker.chunk_text(
                text=extracted_text,
                document_id=str(document.id),
                subject_code=subject_code,
                chapter=chapter,
                faculty_id=faculty_employee_id,
            )

            if not chunks:
                document.status = "ready"
                document.chunk_count = 0
                document.save(
                    update_fields=(
                        "extracted_text",
                        "status",
                        "chunk_count",
                        "error_message",
                    )
                )
                return document

            chunk_ids = [str(uuid.uuid4()) for _ in range(len(chunks))]
            vectors = self._embeddings.embed_texts(
                [c.content for c in chunks]
            )
            payloads = [
                {
                    "content": c.content,
                    "document_id": str(document.id),
                    "subject_code": subject_code,
                    "chapter": chapter,
                    "faculty_id": faculty_employee_id,
                    "chunk_index": c.chunk_index,
                }
                for c in chunks
            ]

            self._vectorstore.upsert_chunks(chunk_ids, vectors, payloads)

            db_chunks = [
                DocumentChunk(
                    document=document,
                    chunk_index=c.chunk_index,
                    content=c.content,
                    page_number=c.page_number,
                    embedding_id=chunk_ids[i],
                    metadata=payloads[i],
                )
                for i, c in enumerate(chunks)
            ]
            DocumentChunk.objects.bulk_create(db_chunks)

            document.chunk_count = len(chunks)
            document.status = "ready"
            document.save(
                update_fields=(
                    "extracted_text",
                    "status",
                    "chunk_count",
                    "error_message",
                )
            )

        except DocumentServiceError:
            raise
        except Exception as exc:
            logger.exception("Document processing failed: %s", exc)
            document.status = "failed"
            document.error_message = str(exc)[:500]
            document.save(update_fields=("status", "error_message"))
            raise DocumentServiceError(
                f"Document processing failed: {exc}"
            ) from exc

        return document

    @transaction.atomic
    def delete_document(self, document: UploadedDocument):
        """Delete document, its chunks, and vector embeddings."""
        self._ensure_services()
        self._vectorstore.delete_by_document(str(document.id))
        document.chunks.all().delete()
        document.file.delete(save=False)
        document.delete()
